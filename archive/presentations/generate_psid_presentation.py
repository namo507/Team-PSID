from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SUMMARY_PATH = ROOT / "psid_artifact_summary.json"
OUTPUT_PATH = ROOT / "PSID_Crisis_Module_Client_Deck.pptx"

FIG_TOP = ROOT / "fig_top_ranked_questions.png"
FIG_SCATTER = ROOT / "fig_utility_vs_burden.png"
FIG_HEATMAP = ROOT / "fig_construct_heatmap.png"
FIG_TIME = ROOT / "fig_time_budget.png"

BG = RGBColor(247, 246, 242)
INK = RGBColor(28, 37, 52)
MUTED = RGBColor(88, 99, 117)
BLUE = RGBColor(37, 99, 235)
RED = RGBColor(157, 23, 77)
ORANGE = RGBColor(217, 119, 6)
GREEN = RGBColor(22, 101, 52)
GOLD = RGBColor(180, 115, 0)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(220, 226, 234)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.75))
    text_frame = box.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.color.rgb = INK

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.88), Inches(12.0), Inches(0.35))
        tf = sub.text_frame
        p2 = tf.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED


def add_footer(slide, text: str):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.28))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(8.5)
    run.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text, *, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_list(slide, left, top, width, height, items, *, size=13, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_kpi_card(slide, left, top, width, height, value, label, accent):
    card = add_panel(slide, left, top, width, height, radius=True)
    card.fill.fore_color.rgb = PANEL
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.07))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    add_text(slide, left + Inches(0.18), top + Inches(0.22), width - Inches(0.3), Inches(0.45), value, size=24, bold=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.72), width - Inches(0.3), Inches(0.45), label, size=11, color=MUTED)


def add_picture_fit(slide, image_path: Path, left, top, width, height):
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def build_slide_1(prs: Presentation, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "PSID Crisis Module: What the Optimized Design Delivers", "Validated portfolio metrics from the final ranked questionnaire workflow")

    add_panel(slide, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.9), radius=True)
    add_text(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(11.6),
        Inches(0.4),
        "The final module keeps the questionnaire inside the 30-minute field constraint while preserving the highest-value, most portable crisis questions.",
        size=16,
        bold=True,
        color=INK,
    )

    cards = [
        (f"{summary['rows']} -> {summary['selected_rows']}", "Candidate questions narrowed to deployable module", BLUE),
        (f"{summary['selected_minutes']:.2f} / 30 min", "Selected module duration under hard interview cap", GREEN),
        (f"{summary['avg_pi']:.3f} / {summary['max_pi']:.3f}", "Average / maximum Pi enhanced-priority score", RED),
        (
            f"{summary['recommended_generic_count'] + summary['recommended_specific_count']} recommendations",
            "Documentation-backed outputs: 6 generic core + 8 crisis-specific",
            ORANGE,
        ),
    ]
    positions = [
        (Inches(0.65), Inches(2.45)),
        (Inches(6.45), Inches(2.45)),
        (Inches(0.65), Inches(4.55)),
        (Inches(6.45), Inches(4.55)),
    ]
    for (value, label, accent), (left, top) in zip(cards, positions):
        add_kpi_card(slide, left, top, Inches(5.6), Inches(1.75), value, label, accent)

    add_bullet_list(
        slide,
        Inches(0.82),
        Inches(6.45),
        Inches(11.2),
        Inches(0.45),
        [
            "Selection outcome: 28 questions retained from 52 ranked candidates, consuming 97.2% of the 30-minute ceiling without exceeding it.",
        ],
        size=12,
    )
    add_footer(slide, "Sources: psid_artifact_summary.json and generate_psid_artifacts.py")


def build_slide_2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Highest-Value Questions and the Utility-Burden Tradeoff", "The charts use the generated PNG artifacts directly to preserve the validated visual output")

    add_panel(slide, Inches(0.45), Inches(1.25), Inches(6.1), Inches(5.45), radius=True)
    add_panel(slide, Inches(6.7), Inches(1.25), Inches(6.1), Inches(5.45), radius=True)
    add_text(slide, Inches(0.7), Inches(1.42), Inches(3.0), Inches(0.3), "Top-ranked questions by baseline utility-to-burden ratio", size=14, bold=True)
    add_text(slide, Inches(6.95), Inches(1.42), Inches(3.5), Inches(0.3), "Utility vs burden profile across the integrated crisis corpus", size=14, bold=True)
    add_picture_fit(slide, FIG_TOP, Inches(0.62), Inches(1.8), Inches(5.75), Inches(4.55))
    add_picture_fit(slide, FIG_SCATTER, Inches(6.86), Inches(1.8), Inches(5.75), Inches(4.55))
    add_bullet_list(
        slide,
        Inches(0.75),
        Inches(6.45),
        Inches(11.5),
        Inches(0.5),
        [
            "Upper-left positions indicate better questions: higher utility and lower respondent burden. Larger bubbles reflect stronger ranking scores.",
        ],
        size=12,
    )
    add_footer(slide, "Embedded visuals: fig_top_ranked_questions.png and fig_utility_vs_burden.png")


def build_slide_3(prs: Presentation, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Coverage Across Constructs, Sources, and Time Budget", "Three complementary views show breadth of coverage, source integration, and budget discipline")

    add_panel(slide, Inches(0.45), Inches(1.2), Inches(7.5), Inches(5.8), radius=True)
    add_panel(slide, Inches(8.1), Inches(1.2), Inches(4.65), Inches(2.75), radius=True)
    add_panel(slide, Inches(8.1), Inches(4.1), Inches(4.65), Inches(2.9), radius=True)

    add_text(slide, Inches(0.68), Inches(1.38), Inches(3.8), Inches(0.3), "Construct coverage by source", size=14, bold=True)
    add_picture_fit(slide, FIG_HEATMAP, Inches(0.62), Inches(1.72), Inches(7.15), Inches(5.0))

    add_text(slide, Inches(8.35), Inches(1.38), Inches(3.0), Inches(0.3), "Selected source contribution", size=14, bold=True)
    source_lines = [
        f"Hurricane Katrina 2007: {summary['selected_source_counts']['Hurricane Katrina 2007']} selected",
        f"COVID-19: {summary['selected_source_counts']['COVID-19']} selected",
        f"Govt Shutdown Income: {summary['selected_source_counts']['Govt Shutdown Income']} selected",
        f"Understanding Society: {summary['selected_source_counts']['Understanding Society']} selected",
        f"Govt Shutdown Crisis: {summary['selected_source_counts']['Govt Shutdown Crisis']} selected",
    ]
    add_bullet_list(slide, Inches(8.35), Inches(1.8), Inches(3.9), Inches(1.7), source_lines, size=12)

    add_text(slide, Inches(8.35), Inches(4.28), Inches(3.3), Inches(0.3), "Time budget by toggle category", size=14, bold=True)
    add_picture_fit(slide, FIG_TIME, Inches(8.23), Inches(4.65), Inches(4.25), Inches(2.02))
    add_footer(slide, "Embedded visuals: fig_construct_heatmap.png and fig_time_budget.png")


def build_slide_4(prs: Presentation, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Metrics That Are Easy to Explain to Stakeholders", "The workflow balances interpretability, deployability, and evidence-backed recommendations")

    add_panel(slide, Inches(0.5), Inches(1.25), Inches(4.0), Inches(5.55), radius=True)
    add_panel(slide, Inches(4.75), Inches(1.25), Inches(4.0), Inches(5.55), radius=True)
    add_panel(slide, Inches(9.0), Inches(1.25), Inches(3.75), Inches(5.55), radius=True)

    add_text(slide, Inches(0.72), Inches(1.48), Inches(2.7), Inches(0.3), "Core performance metrics", size=14, bold=True)
    add_bullet_list(
        slide,
        Inches(0.74),
        Inches(1.88),
        Inches(3.4),
        Inches(4.6),
        [
            f"52 total ranked questions; 28 selected for the deployable module.",
            f"29.17 selected minutes versus 68.95 minutes for the full corpus.",
            f"Average Ri: {summary['avg_ri']:.3f}; maximum Ri: {summary['max_ri']:.3f}.",
            f"Average Pi: {summary['avg_pi']:.3f}; maximum Pi: {summary['max_pi']:.3f}.",
        ],
        size=13,
    )

    add_text(slide, Inches(4.98), Inches(1.48), Inches(2.7), Inches(0.3), "Simple interpretation", size=14, bold=True)
    add_bullet_list(
        slide,
        Inches(5.0),
        Inches(1.88),
        Inches(3.3),
        Inches(4.6),
        [
            "Ri answers a basic question: how much thematic value do we get per unit of burden?",
            "Pi improves that score by rewarding rarer wording, richer construct coverage, and deployable phrasing while penalizing redundancy.",
            "The selected module is near the cap because the workflow fills the available interview time with the highest-priority items first.",
        ],
        size=13,
    )

    add_text(slide, Inches(9.22), Inches(1.48), Inches(2.4), Inches(0.3), "Recommendation outputs", size=14, bold=True)
    add_kpi_card(slide, Inches(9.25), Inches(1.95), Inches(3.05), Inches(1.45), str(summary['recommended_generic_count']), "Generic-core recommendations", BLUE)
    add_kpi_card(slide, Inches(9.25), Inches(3.6), Inches(3.05), Inches(1.45), str(summary['recommended_specific_count']), "Crisis-specific recommendations", RED)
    add_text(
        slide,
        Inches(9.28),
        Inches(5.35),
        Inches(3.0),
        Inches(0.8),
        "Each recommendation is linked to documentation-backed patterns and supported by the highest-scoring matching questions in the corpus.",
        size=11,
        color=MUTED,
    )
    add_footer(slide, "Metrics derived from psid_artifact_summary.json and recommendation generation in generate_psid_artifacts.py")


def add_stage(slide, left, top, width, height, title, lines, accent):
    box = add_panel(slide, left, top, width, height, radius=True)
    box.fill.fore_color.rgb = RGBColor(252, 252, 251)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.12), top + Inches(0.18), width - Inches(0.24), Inches(0.32), title, size=13, bold=True)
    add_bullet_list(slide, left + Inches(0.12), top + Inches(0.56), width - Inches(0.24), height - Inches(0.7), lines, size=10, color=MUTED)
    return box


def build_slide_5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Computational Model Architecture", "From raw PSID source material to deployable, evidence-backed questionnaire artifacts")

    stages = [
        (Inches(0.4), Inches(1.55), Inches(2.15), Inches(2.3), "1. Data Input", ["5 PSID-related sources", "52 integrated candidate questions", "Authoritative base CSV for ranking"], BLUE),
        (Inches(2.8), Inches(1.55), Inches(2.15), Inches(2.3), "2. NLP Processing", ["Keyword parsing and tagging", "Construct extraction", "Question-level feature assembly"], GREEN),
        (Inches(5.2), Inches(1.55), Inches(2.15), Inches(2.3), "3. Augmented Scoring", ["TF-IDF rarity", "Cosine redundancy penalty", "Construct richness and portability bonus"], RED),
        (Inches(7.6), Inches(1.55), Inches(2.15), Inches(2.3), "4. Selection Engine", ["Rank by Pi descending", "Greedy fit to 30-minute cap", "Deployable wording and recommendations"], ORANGE),
        (Inches(10.0), Inches(1.55), Inches(2.4), Inches(2.3), "5. Artifact Output", ["Final CSV with 25 columns", "Dashboard data payload", "PNG figures, report HTML, notebook outputs"], GOLD),
    ]
    shapes = [add_stage(slide, *stage) for stage in stages]
    for idx in range(len(shapes) - 1):
        x1 = shapes[idx].left + shapes[idx].width
        y1 = shapes[idx].top + shapes[idx].height // 2
        x2 = shapes[idx + 1].left
        y2 = shapes[idx + 1].top + shapes[idx + 1].height // 2
        line = slide.shapes.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = MUTED
        line.line.width = Pt(2)

    formula = add_panel(slide, Inches(0.75), Inches(4.35), Inches(7.2), Inches(1.8), radius=True)
    formula.fill.fore_color.rgb = RGBColor(255, 252, 245)
    add_text(slide, Inches(0.98), Inches(4.58), Inches(1.8), Inches(0.3), "Enhanced priority formula", size=14, bold=True)
    add_text(
        slide,
        Inches(0.98),
        Inches(4.98),
        Inches(6.7),
        Inches(0.75),
        "Pi = U* / [Bi x (1 + 0.65 x redundancy)]\nU* = U x (1 + 0.32 x idf + 0.24 x construct + 0.12 x richness + portability)",
        size=13,
        color=INK,
    )

    legend = add_panel(slide, Inches(8.2), Inches(4.35), Inches(4.45), Inches(1.8), radius=True)
    add_text(slide, Inches(8.42), Inches(4.58), Inches(2.3), Inches(0.3), "Construct priority weights", size=14, bold=True)
    add_text(
        slide,
        Inches(8.42),
        Inches(4.98),
        Inches(3.9),
        Inches(0.8),
        "Trauma/Health 0.58\nHousing/Shelter 0.55\nGovernment Aid 0.48\nFinancial Coping 0.45\nEmployment 0.42\nEconomic/Income 0.40\nDemographics 0.14",
        size=11,
        color=MUTED,
    )
    add_footer(slide, "Implemented in generate_psid_artifacts.py via build_all() with TF-IDF, redundancy control, and recommendation generation")


def main():
    summary = json.loads(SUMMARY_PATH.read_text())
    for image in (FIG_TOP, FIG_SCATTER, FIG_HEATMAP, FIG_TIME):
        if not image.exists():
            raise FileNotFoundError(f"Missing required figure: {image}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs, summary)
    build_slide_2(prs)
    build_slide_3(prs, summary)
    build_slide_4(prs, summary)
    build_slide_5(prs)

    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()