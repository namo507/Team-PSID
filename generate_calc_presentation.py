from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SUMMARY_PATH = ROOT / "psid_artifact_summary.json"
OUTPUT_PATH = ROOT / "calc.pptx"

BG = RGBColor(248, 249, 251)
PANEL = RGBColor(255, 255, 255)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
NAVY = RGBColor(30, 58, 95)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
LINE = RGBColor(229, 231, 235)
SOFT_BLUE = RGBColor(248, 251, 255)
SOFT_AMBER = RGBColor(255, 248, 240)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_panel(slide, left, top, width, height, *, fill_color=PANEL, line_color=LINE, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text, *, size=12, bold=False, color=INK, font_name="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, size=11, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
    return box


def add_code(slide, left, top, width, height, text, *, size=11, fill_color=SOFT_BLUE):
    add_panel(slide, left, top, width, height, fill_color=fill_color, rounded=True)
    return add_text(
        slide,
        left + Inches(0.12),
        top + Inches(0.1),
        width - Inches(0.24),
        height - Inches(0.2),
        text,
        size=size,
        color=NAVY,
        font_name="Consolas",
    )


def add_title(slide, title: str, subtitle: str):
    add_text(slide, Inches(0.45), Inches(0.28), Inches(12.2), Inches(0.46), title, size=26, bold=True, color=NAVY, font_name="Aptos Display")
    add_text(slide, Inches(0.47), Inches(0.8), Inches(12.0), Inches(0.32), subtitle, size=11, color=MUTED)


def add_footer(slide, text: str):
    add_text(slide, Inches(0.48), Inches(7.02), Inches(12.2), Inches(0.22), text, size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = BLUE
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True
    return line


def add_metric_card(slide, left, top, width, height, value, label, accent):
    add_panel(slide, left, top, width, height)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.06))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.12), top + Inches(0.16), width - Inches(0.24), Inches(0.36), value, size=20, bold=True)
    add_text(slide, left + Inches(0.12), top + Inches(0.54), width - Inches(0.24), Inches(0.32), label, size=9.5, color=MUTED)


def build_slide_1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "How The Code Produces Ri And Pi", "A 3-stage view of the pipeline from ranked CSV to final stakeholder artifacts")

    stage_width = Inches(1.9)
    stage_height = Inches(1.45)
    stage_top = Inches(1.5)
    stage_lefts = [Inches(0.45), Inches(2.55), Inches(4.65), Inches(6.75), Inches(8.85), Inches(10.95)]
    stages = [
        ("1. Input CSV", ["52 ranked questions", "5 PSID-related sources", "final corpus file"]),
        ("2. load_dataset()", ["read CSV", "derive selected flag", "estimate minutes"]),
        ("3. NLP helpers", ["parse keywords", "tag crisis terms", "extract constructs"]),
        ("4. Baseline score", ["compute Ui", "compute Bi", "rank with Ri = Ui / Bi"]),
        ("5. Enhanced priority", ["TF-IDF rarity", "construct bonus", "Pi with redundancy penalty"]),
        ("6. build_all()", ["write summary JSON", "write dashboard JS", "refresh figures + HTML"]),
    ]

    for index, (title, bullets) in enumerate(stages):
        left = stage_lefts[index]
        add_panel(slide, left, stage_top, stage_width, stage_height)
        add_text(slide, left + Inches(0.1), stage_top + Inches(0.12), stage_width - Inches(0.2), Inches(0.25), title, size=12, bold=True, color=NAVY)
        add_bullets(slide, left + Inches(0.1), stage_top + Inches(0.42), stage_width - Inches(0.18), Inches(0.85), bullets, size=9.5, color=MUTED)
        if index < len(stages) - 1:
            add_arrow(slide, left + stage_width, stage_top + Inches(0.72), stage_lefts[index + 1], stage_top + Inches(0.72))

    add_panel(slide, Inches(0.45), Inches(3.35), Inches(6.15), Inches(2.7), fill_color=SOFT_BLUE)
    add_text(slide, Inches(0.68), Inches(3.55), Inches(2.8), Inches(0.24), "Functions Used In The Calculation Path", size=14, bold=True, color=NAVY)
    add_code(
        slide,
        Inches(0.68),
        Inches(3.92),
        Inches(5.7),
        Inches(1.78),
        "load_dataset()\nparse_keywords()\ntag_keywords()\nextract_constructs()\ncompute_augmented_scores()\nconstruct_bonus()\nportability_bonus()\nbuild_all()",
        size=12,
        fill_color=PANEL,
    )

    add_panel(slide, Inches(6.8), Inches(3.35), Inches(5.95), Inches(2.7), fill_color=SOFT_AMBER)
    add_text(slide, Inches(7.02), Inches(3.55), Inches(3.0), Inches(0.24), "How The Final Numbers Emerge", size=14, bold=True, color=NAVY)
    add_bullets(
        slide,
        Inches(7.02),
        Inches(3.92),
        Inches(5.4),
        Inches(1.78),
        [
            "Ri is the first-pass utility-to-burden ranking signal.",
            "Pi upgrades Ri by adding rarity, construct priority, richness, portability, and redundancy control.",
            "The selected module is then constrained to the 30-minute deployment ceiling.",
            "The same code writes the dashboard, report, recommendation bundle, and figure outputs.",
        ],
        size=11,
        color=INK,
    )

    add_footer(slide, "Based on generate_psid_artifacts.py and PSID_NLP_Crisis_Module_Structure.py")


def build_slide_2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "The Exact Formulas Used In Code", "The coefficients below are hard-coded design parameters that generate Ri and Pi")

    add_panel(slide, Inches(0.45), Inches(1.35), Inches(5.95), Inches(5.4))
    add_panel(slide, Inches(6.65), Inches(1.35), Inches(6.15), Inches(5.4))

    add_text(slide, Inches(0.7), Inches(1.58), Inches(2.8), Inches(0.25), "Baseline ranking", size=15, bold=True, color=NAVY)
    add_code(slide, Inches(0.72), Inches(2.0), Inches(5.35), Inches(0.9), "Ri = Ui / Bi", size=22, fill_color=PANEL)
    add_code(slide, Inches(0.72), Inches(3.0), Inches(5.35), Inches(1.0), "Bi = max(0.10 * word_count + 0.20 * complexity, 0.1)", size=17, fill_color=PANEL)
    add_bullets(
        slide,
        Inches(0.74),
        Inches(4.28),
        Inches(5.2),
        Inches(1.8),
        [
            "Ui is the utility signal built from taxonomy-weighted crisis keywords.",
            "Bi is burden from wording length and structural complexity.",
            "ALPHA = 0.10 and BETA = 0.20 come from the notebook support module.",
        ],
        size=11,
        color=INK,
    )

    add_text(slide, Inches(6.9), Inches(1.58), Inches(3.2), Inches(0.25), "Enhanced priority", size=15, bold=True, color=NAVY)
    add_code(slide, Inches(6.92), Inches(2.0), Inches(5.55), Inches(1.02), "Pi = U* / [Bi * (1 + 0.65 * redundancy_scaled)]", size=15, fill_color=PANEL)
    add_code(slide, Inches(6.92), Inches(3.16), Inches(5.55), Inches(1.1), "U* = Ui * (1 + 0.32 * idf_scaled + 0.24 * construct_scaled + 0.12 * richness_scaled + portability_bonus)", size=12.5, fill_color=PANEL)
    add_bullets(
        slide,
        Inches(6.94),
        Inches(4.48),
        Inches(5.35),
        Inches(1.45),
        [
            "0.32 rewards lexical rarity through TF-IDF.",
            "0.24 rewards construct bonus built from rarity and priority weights.",
            "0.12 rewards construct richness.",
            "0.65 penalizes high redundancy using cosine similarity.",
        ],
        size=10.8,
        color=INK,
    )

    add_panel(slide, Inches(0.45), Inches(6.05), Inches(12.35), Inches(0.72), fill_color=SOFT_BLUE)
    add_text(slide, Inches(0.68), Inches(6.26), Inches(12.0), Inches(0.22), "These are explicitly coded workflow coefficients, not regression-estimated coefficients. They define how the scoring system weights burden, rarity, construct strength, richness, portability, and redundancy.", size=10.5, color=MUTED)
    add_footer(slide, "Key constants: ALPHA, BETA, SECS_PER_WORD, MAX_SECONDS in PSID_NLP_Crisis_Module_Structure.py")


def build_slide_3(prs: Presentation, summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "How The Priority Weights Create The Final Outputs", "The construct-weight dictionary feeds construct_bonus(), which feeds Pi, which drives the final benchmark module")

    add_panel(slide, Inches(0.45), Inches(1.35), Inches(4.05), Inches(5.35))
    add_panel(slide, Inches(4.75), Inches(1.35), Inches(3.85), Inches(5.35), fill_color=SOFT_BLUE)
    add_panel(slide, Inches(8.85), Inches(1.35), Inches(3.95), Inches(5.35))

    add_text(slide, Inches(0.68), Inches(1.58), Inches(3.0), Inches(0.25), "CONSTRUCT_PRIORITY weights", size=15, bold=True, color=NAVY)
    weights = [
        "Trauma / Health = 0.58",
        "Housing / Shelter = 0.55",
        "Government Aid = 0.48",
        "Financial Coping = 0.45",
        "Employment = 0.42",
        "Economic / Income = 0.40",
        "Demographics = 0.14",
    ]
    add_code(slide, Inches(0.68), Inches(1.96), Inches(3.58), Inches(3.0), "CONSTRUCT_PRIORITY = {\n  'Trauma / Health': 0.58,\n  'Housing / Shelter': 0.55,\n  'Government Aid': 0.48,\n  'Financial Coping': 0.45,\n  'Employment': 0.42,\n  'Economic / Income': 0.40,\n  'Demographics': 0.14\n}", size=11.2, fill_color=PANEL)
    add_text(slide, Inches(0.72), Inches(5.22), Inches(3.4), Inches(0.8), "These values are defined once in code and then averaged inside construct_bonus() for each question's construct set.", size=10.5, color=MUTED)

    add_text(slide, Inches(4.98), Inches(1.58), Inches(2.8), Inches(0.25), "construct_bonus() logic", size=15, bold=True, color=NAVY)
    add_code(
        slide,
        Inches(5.0),
        Inches(1.96),
        Inches(3.35),
        Inches(2.28),
        "rarity = mean(1 / freq) * N\npriority = mean(CONSTRUCT_PRIORITY[name])\nrichness = len(set(constructs))\nconstruct_bonus = rarity + priority + 0.08 * richness",
        size=12,
        fill_color=PANEL,
    )
    add_code(
        slide,
        Inches(5.0),
        Inches(4.45),
        Inches(3.35),
        Inches(1.18),
        "portability_bonus\nGeneric Core: 0.16 or 0.03\nToggle items: 0.08 or 0.02",
        size=11,
        fill_color=PANEL,
    )

    add_text(slide, Inches(9.08), Inches(1.58), Inches(2.5), Inches(0.25), "Validated output metrics", size=15, bold=True, color=NAVY)
    metric_specs = [
        (f"{summary['rows']}", "total questions", BLUE),
        (f"{summary['selected_rows']}", "selected", GREEN),
        (f"{summary['selected_minutes']:.2f}", "minutes", AMBER),
        (f"{summary['avg_ri']:.3f}", "Avg Ri", BLUE),
        (f"{summary['avg_pi']:.3f}", "Avg Pi", RED),
        (f"{summary['max_pi']:.3f}", "Max Pi", RED),
    ]
    positions = [
        (Inches(9.08), Inches(1.98)),
        (Inches(10.95), Inches(1.98)),
        (Inches(9.08), Inches(3.18)),
        (Inches(10.95), Inches(3.18)),
        (Inches(9.08), Inches(4.38)),
        (Inches(10.95), Inches(4.38)),
    ]
    for (value, label, accent), (left, top) in zip(metric_specs, positions):
        add_metric_card(slide, left, top, Inches(1.7), Inches(1.0), value, label, accent)

    add_text(slide, Inches(9.08), Inches(5.72), Inches(3.2), Inches(0.4), "These are the published benchmark numbers written into psid_artifact_summary.json and reused by the dashboard and report.", size=10.2, color=MUTED)
    add_footer(slide, "Weights: generate_psid_artifacts.py | Constants: PSID_NLP_Crisis_Module_Structure.py | Results: psid_artifact_summary.json")


def main():
    summary = json.loads(SUMMARY_PATH.read_text())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs, summary)

    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()